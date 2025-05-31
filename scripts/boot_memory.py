import glob, os
import pypdf
import pdfplumber
from PIL import Image
import easyocr
import pytesseract
from docx import Document
import openpyxl
from pptx import Presentation
from openai import OpenAI
from pinecone import Pinecone, ServerlessSpec
import re
from datetime import datetime
import json
import logging

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
INDEX_NAME = "companion-memory"

if not OPENAI_API_KEY or not PINECONE_API_KEY:
    raise ValueError("Missing required API keys")

client = OpenAI(api_key=OPENAI_API_KEY)
pc = Pinecone(api_key=PINECONE_API_KEY)

# Ensure index exists
try:
    if INDEX_NAME not in [i["name"] for i in pc.list_indexes()]:
        pc.create_index(
            name=INDEX_NAME, 
            dimension=1536, 
            metric="cosine",
            spec=ServerlessSpec(cloud="aws", region="us-west-2")
        )
    index = pc.Index(INDEX_NAME)
    logger.info(f"Connected to Pinecone index: {INDEX_NAME}")
except Exception as e:
    logger.error(f"Pinecone connection error: {e}")
    raise

# Initialize OCR reader
try:
    ocr_reader = easyocr.Reader(['en'])
    logger.info("EasyOCR initialized for handwriting recognition")
except Exception as e:
    logger.warning(f"OCR initialization warning: {e}")
    ocr_reader = None

def embed(text):  
    try:
        return client.embeddings.create(
            model="text-embedding-3-small", input=text
        ).data[0].embedding
    except Exception as e:
        logger.error(f"Embedding error: {e}")
        return None

def categorize_content(text, filename):
    """Enhanced categorization with better logic"""
    text_lower = text.lower()
    filename_lower = filename.lower()

    # Priority-based categorization with more specific patterns
    if any(word in filename_lower for word in ['journal', 'diary']) or any(word in text_lower for word in ['dear diary', 'today i', 'feeling', 'reflect']):
        if '2022' in filename_lower or '2022' in text_lower:
            return 'journal_2022'
        elif '2024' in filename_lower or '2024' in text_lower:
            return 'journal_2024'
        elif '2025' in filename_lower or '2025' in text_lower:
            return 'journal_2025'
        else:
            return 'personal_journal'
    elif any(word in filename_lower for word in ['personality', 'profile', 'michael']) and not 'company' in filename_lower:
        return 'personality'
    elif any(word in filename_lower for word in ['health', 'medical', 'superbill', 'records']):
        return 'health'
    elif any(word in filename_lower for word in ['company', 'employee', 'handbook', 'rocket', 'launch']):
        return 'work_rls'
    elif any(word in filename_lower for word in ['pitch', 'serwm', 'brand']):
        return 'projects_serwm'
    elif any(word in filename_lower for word in ['attached', 'body', 'myth', 'steal', 'show', 'emyth']):
        return 'books_read'
    elif any(word in text_lower for word in ['goal', 'objective', 'want to', 'plan to', 'achieve']):
        return 'goals'
    elif any(word in text_lower for word in ['meeting', 'call', 'appointment', 'schedule']):
        return 'meetings'
    elif any(word in text_lower for word in ['project', 'task', 'todo', 'work on']):
        return 'projects'
    elif any(word in text_lower for word in ['idea', 'thought', 'concept', 'brainstorm']):
        return 'ideas'
    else:
        return 'general'

def extract_text_with_ocr(file_path):
    """Extract text from images using OCR for handwriting recognition"""
    try:
        if ocr_reader:
            results = ocr_reader.readtext(file_path)
            text = ' '.join([result[1] for result in results])
            return text
        else:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return text
    except Exception as e:
        logger.error(f"OCR error for {file_path}: {e}")
        return ""

def extract_pdf_text_advanced(pdf_path):
    """Advanced PDF text extraction with multiple fallback methods"""
    text = ""

    logger.info(f"Extracting text from: {pdf_path}")

    # Method 1: Try pdfplumber (better for complex layouts)
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text:
                    text += f"\n--- Page {i+1} ---\n{page_text}\n"
        logger.info(f"PDFPlumber extracted {len(text)} characters from {pdf_path}")
    except Exception as e:
        logger.warning(f"PDFPlumber failed for {pdf_path}: {e}")

    # Method 2: Fallback to pypdf if pdfplumber didn't work well
    if len(text.strip()) < 200:
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                fallback_text = ""
                for i, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            fallback_text += f"\n--- Page {i+1} ---\n{page_text}\n"
                    except Exception as e:
                        logger.warning(f"Error extracting page {i+1}: {e}")
                        continue
                if len(fallback_text) > len(text):
                    text = fallback_text
                    logger.info(f"PyPDF extracted {len(text)} characters from {pdf_path}")
        except Exception as e:
            logger.error(f"PyPDF fallback failed for {pdf_path}: {e}")

    return text.strip()

def extract_docx_text(docx_path):
    """Extract text from Word documents"""
    try:
        doc = Document(docx_path)
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        return '\n'.join(text)
    except Exception as e:
        logger.error(f"Error reading DOCX {docx_path}: {e}")
        return ""

def extract_xlsx_text(xlsx_path):
    """Extract text from Excel files"""
    try:
        workbook = openpyxl.load_workbook(xlsx_path)
        text = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text.append(f"Sheet: {sheet_name}")
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                text.append(" | ".join(row_text))
        return '\n'.join(text)
    except Exception as e:
        logger.error(f"Error reading XLSX {xlsx_path}: {e}")
        return ""

def extract_pptx_text(pptx_path):
    """Extract text from PowerPoint files"""
    try:
        prs = Presentation(pptx_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    except Exception as e:
        logger.error(f"Error reading PPTX {pptx_path}: {e}")
        return ""

def chunk_text(text, max_length=700):
    """Improved text chunking with better sentence boundaries"""
    if len(text) <= max_length:
        return [text]

    # Try to split on sentence boundaries first
    sentences = re.split(r'[.!?]+', text)
    chunks = []
    current_chunk = ""

    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue

        if len(current_chunk) + len(sentence) + 2 <= max_length:
            current_chunk += sentence + ". "
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())
            current_chunk = sentence + ". "

    if current_chunk:
        chunks.append(current_chunk.strip())

    # If we still have chunks that are too long, split them
    final_chunks = []
    for chunk in chunks:
        if len(chunk) <= max_length:
            final_chunks.append(chunk)
        else:
            # Split on paragraphs or whitespace
            words = chunk.split()
            temp_chunk = ""
            for word in words:
                if len(temp_chunk) + len(word) + 1 <= max_length:
                    temp_chunk += word + " "
                else:
                    if temp_chunk:
                        final_chunks.append(temp_chunk.strip())
                    temp_chunk = word + " "
            if temp_chunk:
                final_chunks.append(temp_chunk.strip())

    return final_chunks if final_chunks else [text[:max_length]]

def clean_text(text):
    """Clean and normalize text"""
    # Remove excessive whitespace and normalize
    text = re.sub(r'\s+', ' ', text)
    text = text.strip()
    return text

def clear_existing_vectors():
    """Clear existing vectors to avoid duplicates"""
    try:
        # Get all vector IDs and delete them
        stats = index.describe_index_stats()
        if stats.total_vector_count > 0:
            logger.info(f"Clearing {stats.total_vector_count} existing vectors...")
            index.delete(delete_all=True)
            logger.info("Cleared existing vectors")
    except Exception as e:
        logger.warning(f"Error clearing vectors: {e}")

vector_count = 0
processed_files = []
categories_used = set()

print("Starting comprehensive memory synchronization...")
print(f"Working directory: {os.getcwd()}")
print(f"Looking for files in: {os.path.abspath('docs')}")

# Clear existing vectors to avoid duplicates
clear_existing_vectors()

# Process markdown files
md_files = glob.glob("docs/**/*.md", recursive=True)
print(f"Found {len(md_files)} markdown files: {md_files}")

for path in md_files:
    try:
        with open(path, 'r', encoding='utf-8') as f:
            txt = f.read()

        if not txt.strip():
            continue

        txt = clean_text(txt)
        category = categorize_content(txt, path)
        categories_used.add(category)
        chunks = chunk_text(txt)

        for i, chunk in enumerate(chunks):
            if not chunk.strip():
                continue

            embedding = embed(chunk)
            if embedding is None:
                continue

            vector_id = f"md_{os.path.basename(path)}_{i}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            metadata = {
                "text": chunk,
                "source": path,
                "category": category,
                "chunk_index": i,
                "total_chunks": len(chunks),
                "timestamp": datetime.now().isoformat(),
                "file_type": "markdown"
            }

            index.upsert([(vector_id, embedding, metadata)])
            vector_count += 1

        processed_files.append(path)
        print(f"✓ Processed: {path} ({category}) - {len(chunks)} chunks")

    except Exception as e:
        print(f"✗ Error processing {path}: {e}")

# Process PDF files with explicit paths
pdf_patterns = [
    "docs/*.pdf",
    "docs/**/*.pdf"
]

all_pdf_files = []
for pattern in pdf_patterns:
    all_pdf_files.extend(glob.glob(pattern, recursive=True))

# Remove duplicates
pdf_files = list(set(all_pdf_files))
print(f"Found {len(pdf_files)} PDF files:")
for pdf_file in pdf_files:
    print(f"  - {pdf_file}")

def extract_text_from_pdf(pdf_path):
    """Extract text from PDF with OCR fallback for handwritten content"""
    text = ""

    try:
        # Try pdfplumber first (better for formatted PDFs)
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text and len(page_text.strip()) > 50:
                    text += f"\n--- Page {page.page_number} ---\n{page_text}\n"
                else:
                    # If no text or very little text, try OCR
                    print(f"Low text content on page {page.page_number}, attempting OCR...")
                    try:
                        # Convert page to image for OCR
                        img = page.to_image(resolution=300)
                        pil_img = img.original

                        # Use both pytesseract and easyocr for better results
                        ocr_text = ""

                        # Try pytesseract first
                        try:
                            tesseract_text = pytesseract.image_to_string(pil_img)
                            if tesseract_text.strip():
                                ocr_text += tesseract_text
                        except:
                            pass

                        # Also try easyocr for handwriting
                        try:
                            reader = easyocr.Reader(['en'])
                            img_array = np.array(pil_img)
                            easyocr_results = reader.readtext(img_array)
                            easyocr_text = ' '.join([result[1] for result in easyocr_results])
                            if easyocr_text.strip():
                                ocr_text += f"\n[EasyOCR]: {easyocr_text}"
                        except:
                            pass

                        if ocr_text.strip():
                            text += f"\n--- Page {page.page_number} (OCR) ---\n{ocr_text}\n"
                            print(f"OCR extracted {len(ocr_text)} characters from page {page.page_number}")

                    except Exception as ocr_e:
                        print(f"OCR failed for page {page.page_number}: {ocr_e}")

    except Exception as e:
        print(f"pdfplumber failed for {pdf_path}: {e}")

        # Fallback to pypdf
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n--- Page {page_num + 1} ---\n{page_text}\n"
        except Exception as e2:
            print(f"pypdf also failed for {pdf_path}: {e2}")
            return None

    return text.strip() if text.strip() else None

def extract_personal_details(text, source, category):
    """Extract specific personal details like names, relationships, etc."""
    personal_details = []

    # Look for family member mentions
    family_patterns = [
        r"(?:my |his |her )?(?:mom|mother|mama|ma)\s+(?:is\s+)?(\w+)",
        r"(?:my |his |her )?(?:dad|father|papa|pa)\s+(?:is\s+)?(\w+)",
        r"(?:my |his |her )?(?:wife|husband|spouse)\s+(?:is\s+)?(\w+)",
        r"(?:my |his |her )?(?:sister|brother|sibling)\s+(?:is\s+)?(\w+)",
    ]

    import re
    for pattern in family_patterns:
        matches = re.finditer(pattern, text, re.IGNORECASE)
        for match in matches:
            relationship = pattern.split('?')[2].split(')')[0].strip('(').split('|')[0]
            name = match.group(1)
            if name and len(name) > 1 and name.isalpha():
                detail = f"Family member - {relationship}: {name}"
                personal_details.append({
                    'detail': detail,
                    'source': source,
                    'category': f'{category}_family',
                    'confidence': 'high' if category in ['personality', 'journal_2025', 'journal_2024', 'journal_2022'] else 'medium'
                })

    return personal_details

for path in pdf_files:
    try:
        print(f"Processing PDF: {path}")
        txt = extract_text_from_pdf(path)

        if not txt.strip():
            print(f"❌ No text extracted from {path}")
            continue

        print(f"📄 Extracted {len(txt)} characters from {path}")
        txt = clean_text(txt)
        category = categorize_content(txt, path)
        categories_used.add(category)
        chunks = chunk_text(txt, max_length=800)

        processed_chunks = 0
        for i, chunk in enumerate(chunks):
            if not chunk.strip():
                continue

            embedding = embed(chunk)
            if embedding is None:
                continue

            vector_id = f"pdf_{os.path.basename(path)}_{i}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            metadata = {
                "text": chunk,
                "source": path,
                "category": category,
                "chunk_index": i,
                "total_chunks": len(chunks),
                "timestamp": datetime.now().isoformat(),
                "file_type": "pdf"
            }

            index.upsert([(vector_id, embedding, metadata)])
            vector_count += 1
            processed_chunks += 1

        processed_files.append(path)
        print(f"✓ Processed: {path} ({category}) - {processed_chunks} chunks")

    except Exception as e:
        print(f"✗ Error processing PDF {path}: {e}")

# Process Office documents
office_extensions = [
    ('docs/**/*.docx', extract_docx_text),
    ('docs/**/*.xlsx', extract_xlsx_text),
    ('docs/**/*.pptx', extract_pptx_text)
]

for pattern, extract_func in office_extensions:
    files = glob.glob(pattern, recursive=True)
    for path in files:
        try:
            print(f"Processing Office document: {path}")
            txt = extract_func(path)

            if not txt.strip():
                continue

            txt = clean_text(txt)
            category = categorize_content(txt, path)
            categories_used.add(category)
            chunks = chunk_text(txt)

            for i, chunk in enumerate(chunks):
                if not chunk.strip():
                    continue

                embedding = embed(chunk)
                if embedding is None:
                    continue

                vector_id = f"office_{os.path.basename(path)}_{i}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
                metadata = {
                    "text": chunk,
                    "source": path,
                    "category": category,
                    "chunk_index": i,
                    "total_chunks": len(chunks),
                    "timestamp": datetime.now().isoformat(),
                    "file_type": path.split('.')[-1]
                }

                index.upsert([(vector_id, embedding, metadata)])
                vector_count += 1

            processed_files.append(path)
            print(f"✓ Processed: {path} ({category})")

        except Exception as e:
            print(f"✗ Error processing {path}: {e}")

# Process images for handwriting recognition
image_patterns = [
    'docs/**/*.jpg', 'docs/**/*.jpeg', 'docs/**/*.png', 
    'docs/**/*.tiff', 'docs/**/*.bmp'
]

image_files = []
for pattern in image_patterns:
    image_files.extend(glob.glob(pattern, recursive=True))

if image_files:
    print(f"Found {len(image_files)} image files for OCR processing")

    for path in image_files:
        try:
            print(f"Processing image with OCR: {path}")
            txt = extract_text_with_ocr(path)

            if not txt.strip() or len(txt) < 10:
                continue

            txt = clean_text(txt)
            category = categorize_content(txt, path)
            categories_used.add(category)

            embedding = embed(txt)
            if embedding is None:
                continue

            vector_id = f"ocr_{os.path.basename(path)}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            metadata = {
                "text": txt,
                "source": path,
                "category": category,
                "timestamp": datetime.now().isoformat(),
                "file_type": "image_ocr"
            }

            index.upsert([(vector_id, embedding, metadata)])
            vector_count += 1
            processed_files.append(path)
            print(f"✓ Processed OCR: {path} ({category})")

        except Exception as e:
            print(f"✗ Error processing image {path}: {e}")

total_files = len(processed_files)
print(f"\n🎉 Memory sync complete!")
print(f"📁 Successfully processed {total_files} files")
print(f"🧠 Created {vector_count} memory vectors")
print(f"📂 Categories used: {sorted(categories_used)}")
print(f"📊 Average vectors per file: {vector_count/total_files if total_files > 0 else 0:.1f}")

# Verify the index has data
try:
    stats = index.describe_index_stats()
    print(f"📈 Total vectors in index: {stats.total_vector_count}")
except Exception as e:
    print(f"Error getting index stats: {e}")